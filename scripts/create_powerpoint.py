import os
from pptx import Presentation

def main():
    markdown_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'presentation.md')
    output_path = os.path.join(os.path.dirname(__file__), '..', 'docs', 'presentation.pptx')

    try:
        with open(markdown_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except FileNotFoundError:
        print(f"Error: Markdown file not found at {markdown_path}")
        return

    prs = Presentation()
    
    # Split markdown into major sections separated by '---'
    sections = content.split('\n---\n')

    # First section is the title slide
    if sections:
        title_section_lines = sections[0].strip().split('\n')
        title_text = ""
        subtitle_text = []
        if title_section_lines:
            title_text = title_section_lines[0].lstrip('# ').strip()
            for line in title_section_lines[1:]:
                subtitle_text.append(line.lstrip('# ').strip())
        
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        if slide.placeholders[1] and hasattr(slide.placeholders[1], 'text_frame'):
            slide.placeholders[1].text_frame.text = '\n'.join(subtitle_text)

    # Process remaining sections
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    for section in sections[1:]:
        lines = section.strip().split('\n')
        
        current_slide_title = ""
        current_slide_content = []

        for line in lines:
            is_title = False
            if line.startswith('## ') or line.startswith('### '):
                is_title = True

            if is_title and current_slide_title:
                # Add the completed slide
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = current_slide_title.lstrip('#* ').rstrip('* ').strip()
                body_shape = slide.placeholders[1]
                if body_shape and hasattr(body_shape, 'text_frame'):
                    tf = body_shape.text_frame
                    tf.text = '\n'.join(current_slide_content).strip()
                
                # Start a new slide
                current_slide_title = line
                current_slide_content = []
            elif is_title and not current_slide_title:
                current_slide_title = line
            else:
                # Ignore empty lines between titles
                if line.strip() or current_slide_content:
                    current_slide_content.append(line)
        
        # Add the last slide in the section
        if current_slide_title:
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = current_slide_title.lstrip('#* ').rstrip('* ').strip()
            body_shape = slide.placeholders[1]
            if body_shape and hasattr(body_shape, 'text_frame'):
                tf = body_shape.text_frame
                tf.text = '\n'.join(current_slide_content).strip()

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    main()
